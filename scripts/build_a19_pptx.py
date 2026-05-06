"""Generate the A-19 IT Chairperson presentation as A4 portrait .pptx.

Renders each page of the PDF deck to a high-res PNG and embeds it as a
full-bleed slide image. Slide size is set to A4 portrait so the file
opens at the same physical size as the PDF.

Run:  python3 scripts/build_a19_pptx.py
Output: docs/A19_IT_CHAIRPERSON_DECK.pptx
"""
from __future__ import annotations

from pathlib import Path

import pypdfium2 as pdfium
from pptx import Presentation
from pptx.util import Cm

ROOT = Path(__file__).resolve().parent.parent
PDF = ROOT / "docs" / "A19_IT_CHAIRPERSON_DECK.pdf"
OUT = ROOT / "docs" / "A19_IT_CHAIRPERSON_DECK.pptx"
RENDER_DIR = ROOT / "docs" / "_pages"


def render_pages(pdf_path: Path, out_dir: Path, scale: float = 3.0) -> list[Path]:
    out_dir.mkdir(parents=True, exist_ok=True)
    doc = pdfium.PdfDocument(str(pdf_path))
    paths: list[Path] = []
    for i in range(len(doc)):
        pil = doc[i].render(scale=scale).to_pil()
        p = out_dir / f"page_{i + 1:02d}.png"
        pil.save(p, "PNG")
        paths.append(p)
    return paths


def build() -> Path:
    pages = render_pages(PDF, RENDER_DIR, scale=3.0)

    prs = Presentation()
    # A4 portrait
    prs.slide_width = Cm(21.0)
    prs.slide_height = Cm(29.7)

    blank_layout = prs.slide_layouts[6]  # blank
    for img in pages:
        slide = prs.slides.add_slide(blank_layout)
        slide.shapes.add_picture(
            str(img), 0, 0,
            width=prs.slide_width, height=prs.slide_height,
        )

    prs.save(str(OUT))
    return OUT


if __name__ == "__main__":
    path = build()
    print(f"Wrote {path}")
